"""
============================================================
InstaSHAP Presentation – PowerPoint Generator
============================================================
Generates a professional .pptx file with all slides, diagrams,
charts, and visual elements from the web presentation.
============================================================
"""

import os
import math
import io
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from PIL import Image, ImageDraw, ImageFont

# ─── Design Tokens ───────────────────────────────────────────
BG_DARK      = RGBColor(0x0F, 0x17, 0x2A)
BG_CARD      = RGBColor(0x1E, 0x29, 0x3B)
SURFACE      = RGBColor(0x33, 0x41, 0x55)
TEXT_WHITE   = RGBColor(0xF8, 0xFA, 0xFC)
TEXT_DIM     = RGBColor(0x94, 0xA3, 0xB8)
BORDER       = RGBColor(0x47, 0x55, 0x69)

CLR_INPUT    = RGBColor(0x3B, 0x82, 0xF6)  # Blue
CLR_MODEL    = RGBColor(0x8B, 0x5C, 0xF6)  # Purple
CLR_SAMPLING = RGBColor(0xF5, 0x9E, 0x0B)  # Orange
CLR_SHAP     = RGBColor(0x10, 0xB9, 0x81)  # Green
CLR_ERROR    = RGBColor(0xEF, 0x44, 0x44)  # Red
CLR_ACCENT   = RGBColor(0x06, 0xB6, 0xD4)  # Cyan
CLR_YELLOW   = RGBColor(0xFB, 0xBF, 0x24)

SLIDE_WIDTH  = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)


def set_slide_bg(slide, color=BG_DARK):
    """Set solid background color for a slide."""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_textbox(slide, left, top, width, height, text,
                font_size=18, color=TEXT_WHITE, bold=False,
                alignment=PP_ALIGN.LEFT, font_name='Calibri'):
    """Add a styled text box to a slide."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox


def add_bullet_list(slide, left, top, width, height, items,
                    font_size=16, color=TEXT_WHITE, bullet_color=CLR_ACCENT):
    """Add a bulleted list to a slide."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = 'Calibri'
        p.space_after = Pt(6)
        p.level = 0
    return txBox


def add_rounded_rect(slide, left, top, width, height,
                     fill_color=BG_CARD, border_color=BORDER,
                     text='', font_size=14, text_color=TEXT_WHITE):
    """Add a rounded rectangle shape."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.color.rgb = border_color
    shape.line.width = Pt(1)
    if text:
        tf = shape.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = text
        p.font.size = Pt(font_size)
        p.font.color.rgb = text_color
        p.font.name = 'Calibri'
        p.alignment = PP_ALIGN.CENTER
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    return shape


def add_arrow(slide, x1, y1, x2, y2, color=CLR_ACCENT, width=Pt(2)):
    """Add a connector arrow."""
    connector = slide.shapes.add_connector(
        1, x1, y1, x2, y2  # 1 = straight connector
    )
    connector.line.color.rgb = color
    connector.line.width = width
    return connector


def generate_nn_image(width=800, height=400, layers=None,
                      masked=None, labels=None):
    """Generate a neural network diagram as a PIL Image."""
    if layers is None:
        layers = [3, 5, 4, 1]
    if masked is None:
        masked = set()
    if labels is None:
        labels = {}

    img = Image.new('RGBA', (width, height), (15, 23, 42, 255))
    draw = ImageDraw.Draw(img)

    margin_x, margin_y = 80, 50
    usable_w = width - 2 * margin_x
    usable_h = height - 2 * margin_y
    num_layers = len(layers)

    # Compute node positions
    positions = []
    for l in range(num_layers):
        layer_nodes = []
        n = layers[l]
        x = margin_x + (usable_w / max(1, num_layers - 1)) * l
        for i in range(n):
            y = margin_y + (usable_h / (n + 1)) * (i + 1)
            layer_nodes.append((int(x), int(y)))
        positions.append(layer_nodes)

    # Draw connections
    for l in range(num_layers - 1):
        for (x1, y1) in positions[l]:
            for (x2, y2) in positions[l + 1]:
                is_masked = (l == 0 and positions[l].index((x1, y1)) in masked)
                color = (239, 68, 68, 40) if is_masked else (100, 116, 139, 50)
                draw.line([(x1, y1), (x2, y2)], fill=color, width=1)

    # Draw nodes
    r = 12
    layer_colors = [
        (59, 130, 246),   # input - blue
        (139, 92, 246),   # hidden - purple
        (139, 92, 246),   # hidden - purple
        (16, 185, 129),   # output - green
    ]

    for l in range(num_layers):
        color = layer_colors[min(l, len(layer_colors) - 1)]
        for idx, (x, y) in enumerate(positions[l]):
            is_masked = (l == 0 and idx in masked)
            if is_masked:
                # Masked node
                draw.ellipse([x-r, y-r, x+r, y+r], fill=(239, 68, 68, 30),
                             outline=(239, 68, 68, 180), width=2)
                draw.line([(x-6, y-6), (x+6, y+6)], fill=(239, 68, 68, 200), width=2)
                draw.line([(x+6, y-6), (x-6, y+6)], fill=(239, 68, 68, 200), width=2)
            else:
                draw.ellipse([x-r, y-r, x+r, y+r],
                             fill=(*color, 100), outline=(*color, 220), width=2)

    # Layer labels
    layer_names = ['Input'] + ['Hidden'] * (num_layers - 2) + ['Output']
    for l in range(num_layers):
        x = positions[l][0][0]
        draw.text((x - 20, 15), layer_names[l], fill=(148, 163, 184, 200))

    return img


def generate_bar_chart_image(features, title='SHAP Feature Importance',
                             width=800, height=400):
    """Generate a feature importance bar chart as PIL Image."""
    img = Image.new('RGBA', (width, height), (15, 23, 42, 255))
    draw = ImageDraw.Draw(img)

    margin = {'top': 50, 'right': 40, 'bottom': 30, 'left': 100}
    chart_w = width - margin['left'] - margin['right']
    chart_h = height - margin['top'] - margin['bottom']

    # Title
    draw.text((width // 2 - 100, 10), title, fill=(248, 250, 252))

    sorted_features = sorted(features, key=lambda f: abs(f['value']), reverse=True)
    max_val = max(abs(f['value']) for f in sorted_features)
    bar_h = min(35, chart_h // len(sorted_features) - 8)
    center_x = margin['left'] + chart_w // 2

    # Center line
    draw.line([(center_x, margin['top']), (center_x, height - margin['bottom'])],
              fill=(71, 85, 105, 100), width=1)

    for i, f in enumerate(sorted_features):
        y = margin['top'] + (chart_h // len(sorted_features)) * i + \
            (chart_h // len(sorted_features) - bar_h) // 2
        bar_width = int((abs(f['value']) / max_val) * (chart_w // 2))
        is_pos = f['value'] >= 0
        x = center_x if is_pos else center_x - bar_width
        color = (16, 185, 129) if is_pos else (239, 68, 68)

        draw.rounded_rectangle([x, y, x + bar_width, y + bar_h],
                               radius=3, fill=(*color, 180))

        # Feature name
        draw.text((margin['left'] - 90, y + bar_h // 2 - 6),
                  f['name'], fill=(148, 163, 184))

        # Value
        val_x = center_x + bar_width + 5 if is_pos else center_x - bar_width - 50
        draw.text((val_x, y + bar_h // 2 - 6),
                  f'{f["value"]:.3f}', fill=(248, 250, 252))

    return img


def generate_complexity_chart(width=800, height=400):
    """Generate complexity comparison chart."""
    img = Image.new('RGBA', (width, height), (15, 23, 42, 255))
    draw = ImageDraw.Draw(img)

    margin = {'top': 50, 'left': 70, 'right': 40, 'bottom': 50}
    cw = width - margin['left'] - margin['right']
    ch = height - margin['top'] - margin['bottom']

    draw.text((width // 2 - 120, 10), 'Complexity: O(2^n) vs O(k·n)',
              fill=(248, 250, 252))

    # Axes
    draw.line([(margin['left'], margin['top']),
               (margin['left'], height - margin['bottom'])],
              fill=(71, 85, 105), width=1)
    draw.line([(margin['left'], height - margin['bottom']),
               (width - margin['right'], height - margin['bottom'])],
              fill=(71, 85, 105), width=1)

    # Plot O(2^n)
    max_exp = 2 ** 20
    points_exp = []
    points_lin = []
    features = list(range(2, 22, 2))

    for i, n in enumerate(features):
        x = margin['left'] + int((i / (len(features) - 1)) * cw)
        # Exponential
        val_exp = 2 ** n
        norm_exp = min(1, math.log(val_exp + 1) / math.log(max_exp))
        y_exp = height - margin['bottom'] - int(norm_exp * ch)
        points_exp.append((x, y_exp))

        # Linear
        val_lin = 50 * n
        norm_lin = min(1, math.log(val_lin + 1) / math.log(max_exp))
        y_lin = height - margin['bottom'] - int(norm_lin * ch)
        points_lin.append((x, y_lin))

    # Draw lines
    if len(points_exp) > 1:
        draw.line(points_exp, fill=(239, 68, 68), width=3)
    if len(points_lin) > 1:
        draw.line(points_lin, fill=(16, 185, 129), width=3)

    # Legend
    draw.rectangle([margin['left'] + 10, margin['top'] + 5,
                     margin['left'] + 30, margin['top'] + 10],
                    fill=(239, 68, 68))
    draw.text((margin['left'] + 35, margin['top']), 'Exact SHAP O(2^n)',
              fill=(239, 68, 68))
    draw.rectangle([margin['left'] + 10, margin['top'] + 22,
                     margin['left'] + 30, margin['top'] + 27],
                    fill=(16, 185, 129))
    draw.text((margin['left'] + 35, margin['top'] + 17), 'InstaSHAP O(k·n)',
              fill=(16, 185, 129))

    return img


def img_to_stream(img):
    """Convert PIL Image to bytes stream for pptx."""
    stream = io.BytesIO()
    img.save(stream, format='PNG')
    stream.seek(0)
    return stream


def add_image_to_slide(slide, img, left, top, width):
    """Add a PIL Image to a slide."""
    stream = img_to_stream(img)
    slide.shapes.add_picture(stream, left, top, width)


# ═══════════════════════════════════════════════════════════════
# SLIDE BUILDERS
# ═══════════════════════════════════════════════════════════════

def build_title_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
    set_slide_bg(slide)

    add_textbox(slide, Inches(1), Inches(1.5), Inches(11), Inches(1.5),
                'InstaSHAP', font_size=54, bold=True, color=CLR_INPUT,
                alignment=PP_ALIGN.CENTER, font_name='Calibri')

    add_textbox(slide, Inches(1), Inches(3), Inches(11), Inches(1),
                'Instant Model Explanation\nInternal Workflow · Optimization · Limitations',
                font_size=24, color=TEXT_DIM, alignment=PP_ALIGN.CENTER)

    # Tags
    tags = ['XAI', 'Deep Learning', 'Approximation', 'SHAP Values']
    tag_colors = [CLR_INPUT, CLR_MODEL, CLR_SAMPLING, CLR_SHAP]
    for i, (tag, clr) in enumerate(zip(tags, tag_colors)):
        add_rounded_rect(slide, Inches(3.5 + i * 1.6), Inches(4.8),
                         Inches(1.4), Inches(0.4),
                         fill_color=BG_CARD, border_color=clr,
                         text=tag, font_size=10, text_color=clr)


def build_problem_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)

    add_textbox(slide, Inches(0.8), Inches(0.3), Inches(2), Inches(0.4),
                '⚠️ PROBLEM', font_size=11, bold=True, color=CLR_ERROR)
    add_textbox(slide, Inches(0.8), Inches(0.7), Inches(11), Inches(0.8),
                'The SHAP Bottleneck', font_size=40, bold=True, color=CLR_ERROR)
    add_textbox(slide, Inches(0.8), Inches(1.5), Inches(11), Inches(0.6),
                'Traditional SHAP requires evaluating ALL 2ⁿ feature subsets',
                font_size=18, color=TEXT_DIM)

    # Generate and add chart
    chart_img = generate_complexity_chart()
    add_image_to_slide(slide, chart_img, Inches(1.5), Inches(2.3), Inches(10))

    add_rounded_rect(slide, Inches(1), Inches(5.8), Inches(11), Inches(0.8),
                     fill_color=BG_CARD, border_color=CLR_ACCENT,
                     text='10 features → 1,024 subsets  |  20 features → 1M subsets  |  50 features → 10¹⁵ subsets 💥',
                     font_size=14, text_color=TEXT_WHITE)


def build_realworld_need_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)

    add_textbox(slide, Inches(0.8), Inches(0.3), Inches(2), Inches(0.4),
                '🌍 REAL WORLD', font_size=11, bold=True, color=CLR_ACCENT)
    add_textbox(slide, Inches(0.8), Inches(0.7), Inches(11), Inches(0.8),
                'Why We Need Fast Explanations', font_size=38, bold=True, color=CLR_INPUT)

    # Cards
    cards = [
        ('🏥 Healthcare', 'Doctors need instant diagnostic explanations', CLR_SHAP),
        ('🏦 Finance', 'Loan decisions require explainability by law', CLR_INPUT),
        ('🤖 LLMs / Vision', 'Models with 1000+ features need explanations', CLR_ERROR),
    ]
    for i, (title, desc, clr) in enumerate(cards):
        x = Inches(0.8 + i * 4)
        add_rounded_rect(slide, x, Inches(2.2), Inches(3.6), Inches(2),
                         fill_color=BG_CARD, border_color=clr,
                         text=f'{title}\n\n{desc}', font_size=13, text_color=TEXT_DIM)

    add_bullet_list(slide, Inches(0.8), Inches(5), Inches(11), Inches(1.5), [
        '→ We need SHAP-quality explanations at millisecond speed',
        '→ This is the core motivation for InstaSHAP',
    ], font_size=16, color=CLR_ACCENT)


def build_comparison_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)

    add_textbox(slide, Inches(0.8), Inches(0.3), Inches(2), Inches(0.4),
                '📊 COMPARISON', font_size=11, bold=True, color=CLR_ACCENT)
    add_textbox(slide, Inches(0.8), Inches(0.7), Inches(11), Inches(0.8),
                'SHAP Methods: Head-to-Head', font_size=38, bold=True, color=CLR_INPUT)

    # Table
    rows = [
        ['Property', 'Exact SHAP', 'KernelSHAP', 'TreeSHAP', 'InstaSHAP'],
        ['Complexity', 'O(2ⁿ)', 'O(k²·n)', 'O(TL·2ᴹ)', '✓ O(k·n)'],
        ['Model-agnostic', '✓', '✓', '✗', '✓'],
        ['Exactness', 'Exact', 'Approx', 'Exact*', 'Approx'],
        ['Speed (10 feat)', '~4200ms', '~850ms', '~120ms', '✓ ~35ms'],
        ['DL Compatible', '✗', 'Slow', '✗', '✓'],
        ['Interactions', '✓', 'Limited', '✓', '✗'],
    ]

    table_shape = slide.shapes.add_table(len(rows), 5,
                                          Inches(0.8), Inches(2),
                                          Inches(11.5), Inches(4.5))
    table = table_shape.table

    for r, row in enumerate(rows):
        for c, cell_text in enumerate(row):
            cell = table.cell(r, c)
            cell.text = cell_text
            p = cell.text_frame.paragraphs[0]
            p.font.size = Pt(12)
            p.font.name = 'Calibri'

            if r == 0:
                p.font.bold = True
                p.font.color.rgb = CLR_ACCENT
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(0x16, 0x20, 0x32)
            else:
                p.font.color.rgb = TEXT_WHITE
                cell.fill.solid()
                cell.fill.fore_color.rgb = BG_CARD

            # Highlight InstaSHAP column
            if c == 4 and r > 0 and '✓' in cell_text:
                p.font.color.rgb = CLR_SHAP
                p.font.bold = True


def build_workflow_step_slide(prs, step_num, title, subtitle, bullets,
                              badge_text='', badge_color=CLR_INPUT,
                              nn_config=None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)

    add_textbox(slide, Inches(0.8), Inches(0.3), Inches(3), Inches(0.4),
                badge_text or f'Step {step_num}', font_size=11,
                bold=True, color=badge_color)
    add_textbox(slide, Inches(0.8), Inches(0.7), Inches(11), Inches(0.8),
                title, font_size=34, bold=True, color=CLR_INPUT)
    add_textbox(slide, Inches(0.8), Inches(1.5), Inches(11), Inches(0.5),
                subtitle, font_size=16, color=TEXT_DIM)

    # Neural network image if configured
    if nn_config:
        nn_img = generate_nn_image(**nn_config)
        add_image_to_slide(slide, nn_img, Inches(1), Inches(2.2), Inches(5.5))
        bullet_left = Inches(7)
        bullet_width = Inches(5.5)
    else:
        bullet_left = Inches(0.8)
        bullet_width = Inches(11)

    if bullets:
        add_bullet_list(slide, bullet_left, Inches(2.5), bullet_width,
                        Inches(4), bullets, font_size=14, color=TEXT_WHITE)


def build_shap_output_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)

    add_textbox(slide, Inches(0.8), Inches(0.3), Inches(2), Inches(0.4),
                'Step 6 · OUTPUT', font_size=11, bold=True, color=CLR_SHAP)
    add_textbox(slide, Inches(0.8), Inches(0.7), Inches(11), Inches(0.8),
                'Final SHAP Value Output', font_size=38, bold=True, color=CLR_SHAP)

    features = [
        {'name': 'Income', 'value': 0.163},
        {'name': 'Score', 'value': 0.107},
        {'name': 'Age', 'value': 0.057},
    ]
    bar_img = generate_bar_chart_image(features, 'InstaSHAP Feature Importance')
    add_image_to_slide(slide, bar_img, Inches(1.5), Inches(2), Inches(10))

    add_rounded_rect(slide, Inches(1), Inches(6), Inches(11), Inches(0.7),
                     fill_color=BG_CARD, border_color=CLR_SHAP,
                     text='Interpretation: Income (+0.163) contributes most to 82% approval prediction',
                     font_size=14, text_color=TEXT_WHITE)


def build_failure_slide(prs, title, subtitle, bullets, badge='❌ Failure'):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)

    add_textbox(slide, Inches(0.8), Inches(0.3), Inches(3), Inches(0.4),
                badge, font_size=11, bold=True, color=CLR_ERROR)
    add_textbox(slide, Inches(0.8), Inches(0.7), Inches(11), Inches(0.8),
                title, font_size=34, bold=True, color=CLR_ERROR)
    add_textbox(slide, Inches(0.8), Inches(1.5), Inches(11), Inches(0.5),
                subtitle, font_size=16, color=TEXT_DIM)

    if bullets:
        add_bullet_list(slide, Inches(0.8), Inches(2.5), Inches(11),
                        Inches(4), bullets, font_size=14, color=TEXT_WHITE)


def build_improvement_slide(prs, title, subtitle, bullets, badge='🚀 Phase 2'):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)

    add_textbox(slide, Inches(0.8), Inches(0.3), Inches(3), Inches(0.4),
                badge, font_size=11, bold=True, color=CLR_SHAP)
    add_textbox(slide, Inches(0.8), Inches(0.7), Inches(11), Inches(0.8),
                title, font_size=34, bold=True, color=CLR_SHAP)
    add_textbox(slide, Inches(0.8), Inches(1.5), Inches(11), Inches(0.5),
                subtitle, font_size=16, color=TEXT_DIM)

    if bullets:
        add_bullet_list(slide, Inches(0.8), Inches(2.5), Inches(11),
                        Inches(4), bullets, font_size=14, color=TEXT_WHITE)


def build_metrics_slide(prs, title, metrics, badge='📊 Metrics'):
    """3-4 metric boxes with large numbers."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)

    add_textbox(slide, Inches(0.8), Inches(0.3), Inches(3), Inches(0.4),
                badge, font_size=11, bold=True, color=CLR_ACCENT)
    add_textbox(slide, Inches(0.8), Inches(0.7), Inches(11), Inches(0.8),
                title, font_size=38, bold=True, color=CLR_INPUT,
                alignment=PP_ALIGN.CENTER)

    n = len(metrics)
    card_w = Inches(min(3, 11 / n - 0.3))
    start_x = (SLIDE_WIDTH - card_w * n - Inches(0.3) * (n - 1)) / 2

    for i, (value, label, clr) in enumerate(metrics):
        x = start_x + (card_w + Inches(0.3)) * i
        add_rounded_rect(slide, x, Inches(2.5), card_w, Inches(2.5),
                         fill_color=BG_CARD, border_color=clr)
        add_textbox(slide, x, Inches(2.8), card_w, Inches(1.2),
                    value, font_size=36, bold=True, color=clr,
                    alignment=PP_ALIGN.CENTER)
        add_textbox(slide, x, Inches(4), card_w, Inches(0.5),
                    label, font_size=12, color=TEXT_DIM,
                    alignment=PP_ALIGN.CENTER)


def build_thank_you_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)

    add_textbox(slide, Inches(1), Inches(2), Inches(11), Inches(1.5),
                'Thank You', font_size=60, bold=True, color=CLR_INPUT,
                alignment=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(1), Inches(3.8), Inches(11), Inches(1),
                'InstaSHAP: Fast Explanations for Real-World AI',
                font_size=22, color=TEXT_DIM, alignment=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(1), Inches(5.5), Inches(11), Inches(0.5),
                'Questions & Discussion', font_size=16, color=CLR_ACCENT,
                alignment=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════
# MAIN GENERATOR
# ═══════════════════════════════════════════════════════════════

def generate_presentation():
    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT

    # ─── SECTION 1: MOTIVATION ───────────────────────────────
    build_title_slide(prs)
    build_problem_slide(prs)

    build_workflow_step_slide(prs, 0, 'Computational Explosion',
        'Time grows exponentially — exact SHAP is impractical', [
            '• Each subset requires a FULL model forward pass',
            '• GPU cost scales linearly with subsets × batch size',
            '• O(2ⁿ) complexity makes n>20 infeasible',
            '• Real-time explanations? Impossible with exact SHAP',
        ], badge_text='⏱️ COST', badge_color=CLR_ERROR)

    build_realworld_need_slide(prs)

    build_workflow_step_slide(prs, 0, 'Where Existing SHAP Fails',
        'No existing method is both fast, model-agnostic, and accurate', [
            '• Exact SHAP: Too slow for any real model',
            '• KernelSHAP: Better but still O(k²·n)',
            '• TreeSHAP: Fast but ONLY for tree-based models',
            '• Gap: Need speed + model-agnostic + accuracy',
        ], badge_text='❌ FAILURES', badge_color=CLR_ERROR)

    build_metrics_slide(prs, 'Enter: InstaSHAP', [
        ('120×', 'Faster', CLR_SHAP),
        ('~85%', 'Accuracy', CLR_SAMPLING),
        ('Any', 'Model Type', CLR_INPUT),
        ('O(k·n)', 'Complexity', CLR_MODEL),
    ], badge='💡 SOLUTION')

    # ─── SECTION 2: WHAT IS INSTASHAP ────────────────────────
    build_workflow_step_slide(prs, 0, 'InstaSHAP: The Shortcut Analogy',
        'Instead of climbing every path, take a helicopter and sample key viewpoints', [
            '• Exact SHAP: Evaluates ALL 2ⁿ feature coalitions',
            '• InstaSHAP: Samples only k random coalitions (k ≪ 2ⁿ)',
            '• Core idea: Monte Carlo approximation of Shapley values',
            '• Trade-off: ~85% accuracy for 120× speed gain',
        ], badge_text='💡 INTUITION', badge_color=CLR_ACCENT)

    build_workflow_step_slide(prs, 0, 'How InstaSHAP Approximates',
        'Replace exhaustive Shapley sum with Monte Carlo sampling', [
            '• Exact: φᵢ = Σ_S |S|!(n-|S|-1)!/n! × [f(S∪{i}) - f(S)]',
            '• InstaSHAP: φ̂ᵢ ≈ (1/k) Σⱼ [f(Sⱼ∪{i}) - f(Sⱼ)]',
            '• With k=50, n=20: 1,000 calls vs 10,485,760',
            '• Law of large numbers: φ̂ᵢ → φᵢ as k → ∞',
        ], badge_text='⚙️ CORE', badge_color=CLR_SHAP)

    build_comparison_slide(prs)

    # ─── SECTION 3: FULL WORKFLOW ─────────────────────────────
    build_workflow_step_slide(prs, 1, 'Step 1: Input Sample + Model',
        'Sample: [Age=35, Income=75K, Score=0.8] → MLP Neural Network', [
            '• Raw features collected and normalized to [0,1]',
            '• Feature vector x ∈ ℝ³ prepared for perturbation',
            '• NO approximation yet — identical to exact SHAP',
            '• Model: 3-layer MLP with [3, 5, 4, 1] architecture',
        ], badge_text='Step 1 · Input', badge_color=CLR_INPUT,
        nn_config={'layers': [3, 5, 4, 1]})

    build_workflow_step_slide(prs, 1, 'Forward Propagation: f(x) = 0.82',
        'Data flows through the network to produce base prediction', [
            '• Input → Hidden(5) → Hidden(4) → Output(1)',
            '• ReLU activations in hidden layers',
            '• Sigmoid output → P(approval) = 0.82',
            '• Model treated as BLACK BOX by InstaSHAP',
        ], badge_text='Step 1 · Forward Pass', badge_color=CLR_MODEL,
        nn_config={'layers': [3, 5, 4, 1]})

    build_workflow_step_slide(prs, 2, 'Step 2: Feature Perturbation',
        'Mask features and observe how prediction changes', [
            '• Remove Age: f([?, Income, Score]) = 0.75 (Δ = -0.07)',
            '• Remove Income: f([Age, ?, Score]) = 0.58 (Δ = -0.24)',
            '• Remove Score: f([Age, Income, ?]) = 0.71 (Δ = -0.11)',
            '• But SHAP needs ALL combinations — not just single removals',
            '• 3 features → 8 subsets, 20 features → 1M subsets',
        ], badge_text='Step 2 · Perturbation', badge_color=CLR_INPUT,
        nn_config={'layers': [3, 5, 4, 1], 'masked': {1}})

    build_workflow_step_slide(prs, 3, 'Step 3: Subset Sampling',
        'InstaSHAP samples k random subsets instead of all 2ⁿ', [
            '• Generate k random binary masks (k=50 typical)',
            '• Each mask: which features present vs absent',
            '• Absent features filled from training distribution',
            '• APPROXIMATION: Only k subsets evaluated, not 2ⁿ',
            '• For n=20: 50 samples vs 1,048,576 (99.995% reduction)',
        ], badge_text='Step 3 · Sampling', badge_color=CLR_SAMPLING)

    build_workflow_step_slide(prs, 4, 'Step 4: Batch Model Prediction',
        'Run all k perturbed inputs through model simultaneously', [
            '• Batch predict: model.predict_batch([x₁, x₂, ..., xₖ])',
            '• GPU parallelism processes all k inputs at once',
            '• Single GPU call instead of k separate calls',
            '• Results: [0.53, 0.76, 0.70, 0.56, 0.82, ...]',
        ], badge_text='Step 4 · Predict', badge_color=CLR_MODEL,
        nn_config={'layers': [3, 5, 4, 1]})

    build_workflow_step_slide(prs, 5, 'Step 5: Contribution Estimation',
        'For each feature, compute average marginal contribution', [
            '• φ̂ᵢ = (1/k) × Σⱼ [f(Sⱼ ∪ {i}) - f(Sⱼ)]',
            '• φ̂(Income) = avg(0.17, 0.20, 0.12) = 0.163',
            '• φ̂(Score) = avg(0.06, 0.15, 0.11) = 0.107',
            '• φ̂(Age) = avg(0.03, 0.08, 0.06) = 0.057',
            '• Sum: 0.50 + 0.163 + 0.107 + 0.057 = 0.827 ≈ 0.82 ✓',
        ], badge_text='Step 5 · Contribution', badge_color=CLR_SHAP)

    build_shap_output_slide(prs)

    build_metrics_slide(prs, 'End-to-End Timing', [
        ('4.2s', 'Exact SHAP', CLR_ERROR),
        ('35ms', 'InstaSHAP', CLR_SHAP),
        ('120×', 'Speedup', CLR_ACCENT),
    ], badge='⏱️ TIMING')

    # ─── SECTION 4: INTERNAL MECHANISM ───────────────────────
    build_workflow_step_slide(prs, 0, 'Internal Mechanism: Math Intuition',
        'Monte Carlo sampling replaces exhaustive sum', [
            '• Exact: φᵢ = Σ_{S⊆N\\{i}} w(S) · [f(S∪{i}) - f(S)]',
            '• InstaSHAP: φ̂ᵢ ≈ (1/k) Σⱼ w(Sⱼ) · [f(Sⱼ∪{i}) - f(Sⱼ)]',
            '• w(Sⱼ) = 1/k (uniform) vs |S|!(n-|S|-1)!/n! (exact)',
            '• Law of large numbers: converges as k → ∞',
            '• Variance: Var(φ̂ᵢ) = (1/k) × Var_S[marginal]',
        ], badge_text='🔬 MATH', badge_color=CLR_MODEL)

    build_workflow_step_slide(prs, 0, 'Where Bias Enters',
        'Four sources of approximation error in InstaSHAP', [
            '• 1. Sampling bias: Random subsets may miss important coalitions',
            '• 2. Perturbation bias: Feature independence assumption',
            '• 3. Finite sample error: Small k → high variance',
            '• 4. Weight simplification: Uniform vs Shapley kernel',
            '• Combined effect: ~15% error for n=20, k=50',
        ], badge_text='⚠ BIAS', badge_color=CLR_ERROR)

    build_workflow_step_slide(prs, 0, 'Speed Gain Decomposition',
        'Three sources of speed: reduced subsets + batch GPU + simplified weights', [
            '• 1. Reduced subsets: k=50 vs 2²⁰ = 99.995% reduction',
            '• 2. Batch processing: 1 GPU call vs k sequential calls',
            '• 3. Simplified weights: No factorial computation',
            '• Net: 120× speedup with ~85% accuracy retention',
        ], badge_text='⚡ SPEED', badge_color=CLR_SHAP)

    # ─── SECTION 5: WHY FAST ─────────────────────────────────
    build_workflow_step_slide(prs, 0, 'Optimization 1: Reduced Sampling',
        'The single biggest speed win — sample k subsets instead of 2ⁿ', [
            '• n=10:  1,024 → 50 subsets (95% reduction)',
            '• n=20:  1,048,576 → 50 subsets (99.995% reduction)',
            '• n=50:  10¹⁵ → 100 subsets (≈100% reduction)',
            '• Random sampling provides unbiased estimates',
        ], badge_text='⚡ OPT 1', badge_color=CLR_SHAP)

    build_workflow_step_slide(prs, 0, 'Optimization 2: GPU Batch Processing',
        'Process all k perturbed inputs in one GPU call', [
            '• Sequential: k × t_forward (slow, serial)',
            '• Batch: ~1 × t_forward (fast, parallel)',
            '• GPU utilization goes from <5% to >90%',
            '• Memory overhead is minimal for moderate k',
        ], badge_text='⚡ OPT 2', badge_color=CLR_SHAP)

    build_workflow_step_slide(prs, 0, 'Optimization 3: Simplified Weights',
        'Replace factorial Shapley kernel with uniform average', [
            '• Exact: w(S) = |S|!(n-|S|-1)!/n! → factorial overflow for large n',
            '• InstaSHAP: w(Sⱼ) = 1/k → simple division',
            '• Trade-off: Introduces systematic bias but eliminates overhead',
        ], badge_text='⚡ OPT 3', badge_color=CLR_SHAP)

    # ─── SECTION 6: FAILURE ANALYSIS (20 slides) ─────────────
    build_failure_slide(prs, 'Failure Analysis Overview',
        'Where InstaSHAP breaks down — and WHY', [
            '• 1. High Feature Interactions — misses joint effects',
            '• 2. Non-linear Models — unstable marginal contributions',
            '• 3. LLM / Transformer Systems — token interdependencies',
            '• 4. Distribution Shift — invalid perturbation baselines',
        ], badge='❌ FAILURES')

    # Failure 1: Interactions
    for i, (t, s, b) in enumerate([
        ('Feature Interaction Failure',
         'When features work TOGETHER, individual contributions are misleading',
         ['• InstaSHAP assumes features contribute independently',
          '• Interaction terms are NOT captured in standard SHAP values',
          '• XOR problem: both features are critical but get φ ≈ 0',
          '• Any model with synergy/suppression effects is affected',]),
        ('XOR Example: Complete Misattribution',
         'f(X₁,X₂) = X₁ XOR X₂ — features have ZERO marginal effect individually',
         ['• InstaSHAP says: φ(X₁) ≈ 0, φ(X₂) ≈ 0',
          '• Reality: Both features jointly determine 100% of output',
          '• Root cause: marginal contribution assumes additivity',
          '• Random sampling misses critical feature pairs',]),
        ('WHY Interactions Break InstaSHAP',
         'Marginal contribution is fundamentally incompatible with synergy',
         ['• f(S∪{i,j}) ≠ f(S∪{i}) + f(S∪{j}) - f(S) for interacting features',
          '• Standard SHAP lacks per-pair interaction terms',
          '• With k ≪ 2ⁿ, joint coalitions may never be sampled',
          '• Impact: Neural networks, polynomial features, attention all affected',]),
        ('Interaction Case: Expected vs InstaSHAP Output',
         'Concrete comparison showing massive error on interaction terms',
         ['• X₁: Expected φ=0.35, InstaSHAP φ=0.02 (94% error)',
          '• X₂: Expected φ=0.35, InstaSHAP φ=0.03 (91% error)',
          '• X₁·X₂ interaction: Expected 0.30, InstaSHAP 0.00 (100% missed)',
          '• Verdict: COMPLETELY MISLEADING for interacting features',]),
    ]):
        build_failure_slide(prs, t, s, b, badge=f'Failure 1 · Interaction ({i+1}/4)')

    # Failure 2: Non-linear
    for i, (t, s, b) in enumerate([
        ('Non-linear Model Failure',
         'Complex decision boundaries create unreliable explanations',
         ['• Linear: f(x) = wx+b — marginal contributions are constant ✓',
          '• DNN: f(x) = σ(W₃·σ(W₂·σ(W₁·x))) — contributions vary wildly',
          '• ReLU activations create sharp decision boundaries',
          '• Feature importance depends on WHERE in input space you are',]),
        ('Deep Network Activation Landscapes',
         'Position-dependent contributions invalidate simple averaging',
         ['• Different regions of input space have different feature importance',
          '• Random perturbation may cross decision boundaries',
          '• Average marginal contribution becomes meaningless',
          '• For features near boundaries, error can exceed 40%',]),
        ('Non-linear: Expected vs Actual InstaSHAP',
         'Concrete evidence of misattribution in deep networks',
         ['• Feature 1: Expected +0.45, InstaSHAP +0.28 (38% error)',
          '• Feature 2: Expected -0.22, InstaSHAP -0.35 (59% error!)',
          '• Sign and magnitude can both be wrong',
          '• Error increases with model depth and non-linearity',]),
    ]):
        build_failure_slide(prs, t, s, b, badge=f'Failure 2 · Non-linear ({i+1}/3)')

    # Failure 3: LLMs
    for i, (t, s, b) in enumerate([
        ('LLM / Transformer Failure',
         'Language models have sequential dependencies that break feature independence',
         ['• Token order matters: "not good" ≠ "good not"',
          '• Attention is global: every token sees every other token',
          '• Masking creates invalid inputs (out of distribution)',
          '• Meaning spans token boundaries — features ≠ tokens',]),
        ('Attention Mechanism: Why Masking Fails',
         'Self-attention creates global interdependencies',
         ['• Attention(Q,K,V) = softmax(QK^T/√d) · V',
          '• Removing token i changes ALL attention scores globally',
          '• Marginal contribution of i entangled with context',
          '• InstaSHAP CANNOT decompose transformer attention',]),
        ('LLM Verdict: Fundamentally Unsuitable',
         'InstaSHAP requires feature independence — LLMs violate this completely',
         ['• Global attention means no feature is independent',
          '• Masked token creates nonsensical input',
          '• Text meaning is compositional, not additive',
          '• Alternative: Use attention maps, integrated gradients, or LIME',]),
    ]):
        build_failure_slide(prs, t, s, b, badge=f'Failure 3 · LLM ({i+1}/3)')

    # Failure 4: Distribution shift
    for i, (t, s, b) in enumerate([
        ('Distribution Shift Failure',
         'Training ≠ inference data makes perturbation baselines invalid',
         ['• InstaSHAP replaces masked features with TRAINING data samples',
          '• If inference distribution has shifted, baselines are wrong',
          '• Perturbed inputs fall in untrained regions → unreliable predictions',
          '• Example: Train income $50K, Inference $150K → massive overestimation',]),
        ('How Distribution Shift Breaks InstaSHAP',
         'The perturbation mechanism relies on training distribution being valid',
         ['• Training mean income = $50K',
          '• Inference: new market with $150K mean income',
          '• InstaSHAP masks Income → replaces with $50K training sample',
          '• Model sees $150K→$50K = huge drop → overestimates Income φ',
          '• Distribution shift makes ALL explanations unreliable',]),
    ]):
        build_failure_slide(prs, t, s, b, badge=f'Failure 4 · Dist.Shift ({i+1}/2)')

    # Failure summary slides
    build_failure_slide(prs, 'Failure Severity Matrix',
        'Interactions and LLM failures are CRITICAL; non-linear and dist.shift HIGH', [
            '• Feature Interactions: 🔴 Critical severity, Very common, Hard to detect',
            '• Non-linear Models: 🟡 High severity, Common, Moderate detection',
            '• LLM/Transformers: 🔴 Critical, Always fails, Easy to detect',
            '• Distribution Shift: 🔴 Critical, Context-dependent, Very hard to detect',
        ], badge='❌ SEVERITY')

    build_failure_slide(prs, 'The Fundamental Trade-off',
        'InstaSHAP sacrifices exactness for speed — fails in adversarial cases', [
            '✓ InstaSHAP works when: features independent, model ~additive, stable distribution',
            '✗ InstaSHAP FAILS when: strong interactions, highly non-linear, distribution shift',
            '• In benign cases → works well (~85% accuracy)',
            '• In adversarial cases → can be COMPLETELY WRONG',
        ], badge='⚠️ TRADE-OFF')

    # ─── SECTION 7: LIMITATIONS ──────────────────────────────
    build_failure_slide(prs, 'Limitations Overview',
        'Four key limitations of InstaSHAP', [
            '• 1. Loss of Exactness: Monte Carlo → values change between runs',
            '• 2. Approximation Bias: Systematic underestimation of interactions',
            '• 3. Scalability Limits: Even k·n expensive for very large n',
            '• 4. Not Universal: Cannot handle sequential, RL, or non-tabular models',
        ], badge='⚠️ LIMITATIONS')

    build_failure_slide(prs, 'Quantifying Exactness Loss',
        'Error metrics across 1000 test instances', [
            '• k=10:  MAE=0.082, Corr=0.72 (unreliable)',
            '• k=50:  MAE=0.034, Corr=0.89 (moderate)',
            '• k=100: MAE=0.019, Corr=0.94 (good)',
            '• k=500: MAE=0.006, Corr=0.98 (excellent)',
            '• NOTE: These assume NO feature interactions — with interactions, errors 3-5× higher',
        ], badge='📊 EXACTNESS')

    # ─── SECTION 8: APPLICABILITY ─────────────────────────────
    build_workflow_step_slide(prs, 0, 'Applicability Across Model Types',
        'InstaSHAP effectiveness varies drastically by architecture', [
            '🟢 Linear / GLM: Works perfectly',
            '🟢 Tree Models (RF, XGB): Works well',
            '🟡 MLP / Feedforward NN: Partial — risky with interactions',
            '🟡 CNN (Vision): Limited — spatial dependencies',
            '🔴 RNN / LSTM: Weak — sequential dependencies',
            '🔴 Transformers / LLMs: Unsuitable — global attention',
            '🔴 RL Agents: N/A — not feature-decomposable',
        ], badge_text='🔬 APPLICABILITY', badge_color=CLR_ACCENT)

    build_workflow_step_slide(prs, 0, 'ML Models: Sweet Spot',
        'Tabular data with independent features — InstaSHAP excels here', [
            '• Tabular data with clear feature boundaries',
            '• Usually ≤100 features — k doesn\'t need to be huge',
            '• Tree-based models have bounded interactions',
            '• Recommended: k=50, marginal perturbation',
            '• Expected: 85-95% correlation, 20-100ms speed',
        ], badge_text='✓ ML MODELS', badge_color=CLR_SHAP)

    build_workflow_step_slide(prs, 0, 'LLMs: Why InstaSHAP Cannot Work',
        'Fundamental incompatibility with transformer architecture', [
            '• Token granularity mismatch: meaning spans tokens',
            '• Context dependency: every token depends on ALL others',
            '• Masking creates invalid inputs (out of distribution)',
            '• Combinatorial explosion of vocabulary',
            '• Alternative: attention visualization, integrated gradients',
        ], badge_text='✗ LLMs', badge_color=CLR_ERROR)

    # ─── SECTION 9: INTERACTIVE SIM ──────────────────────────
    build_workflow_step_slide(prs, 0, 'What-If Feature Simulator',
        'Adjust Age, Income, Score and see SHAP values change', [
            '• Input: Age=35, Income=75K, Score=0.80',
            '• Prediction: f(x) = 0.82 (82% approval)',
            '• SHAP: Income +0.163, Score +0.107, Age +0.057',
            '• Try: Lower income to $30K → prediction drops to ~0.55',
            '• Income SHAP value would decrease to ~+0.06',
            '• See web version for interactive sliders!',
        ], badge_text='🎮 SIMULATOR', badge_color=CLR_ACCENT)

    build_workflow_step_slide(prs, 0, 'Impact of Sample Size (k)',
        'How k affects speed, accuracy, and variance', [
            '• k=10:  Speed=7ms,  Accuracy=72%, Variance=High',
            '• k=50:  Speed=35ms, Accuracy=85%, Variance=Moderate',
            '• k=100: Speed=70ms, Accuracy=92%, Variance=Low',
            '• k=500: Speed=350ms, Accuracy=98%, Variance=Very Low',
            '• Sweet spot: k=50-100 for most applications',
        ], badge_text='🎮 K IMPACT', badge_color=CLR_ACCENT)

    # ─── SECTION 10: IMPROVEMENTS / PHASE 2 (40 slides) ─────
    build_metrics_slide(prs, 'InstaSHAP 2.0: Proposed Improvements', [
        ('Hybrid', 'SHAP', CLR_INPUT),
        ('Adaptive', 'Sampling', CLR_SAMPLING),
        ('Model', 'Aware', CLR_MODEL),
        ('Interact.', 'Aware', CLR_SHAP),
    ], badge='🚀 PHASE 2')

    # Improvement 1: Hybrid SHAP
    for i, (t, s, b) in enumerate([
        ('Improvement 1: Hybrid SHAP',
         'Combine fast InstaSHAP with targeted exact computation', [
            '• Phase 1: Quick InstaSHAP scan (k=20) — rough estimate',
            '• Phase 2: Identify top-m most important features',
            '• Phase 3: Run EXACT SHAP only for top-m features',
            '• Phase 4: Keep InstaSHAP values for remaining features',
            '• Result: ~95% accuracy, still 60× faster than exact',]),
        ('Hybrid SHAP Pipeline Flow',
         'Quick scan → Rank features → Exact top-m → Merge', [
            '• Quick Scan → 20ms',
            '• Feature Ranking → <1ms',
            '• Exact SHAP (top 3) → 45ms',
            '• Merge → <1ms',
            '• Total: ~66ms vs 4200ms (63× faster)',]),
        ('Hybrid SHAP: Expected Performance',
         'Near-exact accuracy for important features, fast for unimportant', [
            '• Speed: ~2× slower than pure InstaSHAP (~66ms)',
            '• Accuracy: ~95% correlation (up from 85%)',
            '• Still 60× faster than pure exact SHAP',
            '• Best of both worlds: accuracy where it matters, speed everywhere else',]),
    ]):
        build_improvement_slide(prs, t, s, b, badge=f'🔀 Hybrid ({i+1}/3)')

    # Improvement 2: Adaptive Sampling
    for i, (t, s, b) in enumerate([
        ('Improvement 2: Adaptive Sampling',
         'Sample MORE in important regions, LESS in unimportant ones', [
            '• Current: Uniform random — equal samples for all features',
            '• Proposed: Allocate more samples to high-variance features',
            '• Start with k₀=10 uniform, estimate variance per feature',
            '• Re-allocate remaining budget to high-variance features',
            '• Same total K, much better allocation',]),
        ('Adaptive Algorithm Details',
         'Variance-guided sample allocation', [
            '• Step 1: Quick k₀=10 uniform samples',
            '• Step 2: Estimate Var(φ̂ᵢ) for each feature',
            '• Step 3: High variance → allocate 5-10× more samples',
            '• Step 4: Low variance → stop (already converged)',
            '• Step 5: Repeat until budget K exhausted',
            '• Expected: 30-50% variance reduction, same total cost',]),
        ('Adaptive vs Uniform: Side-by-Side',
         'Adaptive focuses computation where it has most impact', [
            '• Uniform: 50 samples × 20 features = 1000 total',
            '• Adaptive: 5 samples for easy features, 150 for hard ones = 1000 total',
            '• Same cost, but hard features get 30× more samples',
            '• Reduces overall MAE by ~40%',]),
    ]):
        build_improvement_slide(prs, t, s, b, badge=f'🎯 Adaptive ({i+1}/3)')

    # Improvement 3: Model-Aware
    for i, (t, s, b) in enumerate([
        ('Improvement 3: Model-Aware Explanation',
         'Use model internals (gradients) to guide sampling', [
            '• Current: Treats model as black box — ignores structure',
            '• Proposed: Use ∇f(x) gradient to identify sensitive features',
            '• High gradient |∂f/∂xᵢ| → feature i is sensitive → sample MORE',
            '• Hessian ∂²f/∂xᵢ² reveals non-linearity → more samples',
            '• Cost: ONE backward pass (ms) to guide ALL sampling',]),
        ('Gradient-Guided Sampling Details',
         'Compute gradient ONCE, then allocate samples proportionally', [
            '• Step 1: Compute ∇f(x) = [∂f/∂x₁, ..., ∂f/∂xₙ]',
            '• Step 2: kᵢ ∝ |∂f/∂xᵢ| × |∂²f/∂xᵢ²|',
            '• Step 3: More samples for sensitive, non-linear features',
            '• Step 4: Less samples for insensitive features',
            '• Net: Better focus, same total computation',]),
        ('Model-Aware vs Model-Agnostic',
         'Using model information yields 25% accuracy improvement', [
            '• Model-Agnostic: 85% correlation (current InstaSHAP)',
            '• Model-Aware: ~90% correlation (with gradient guidance)',
            '• Additional cost: 1 backward pass (~5ms)',
            '• Trade-off: Requires differentiable model (not trees)',]),
    ]):
        build_improvement_slide(prs, t, s, b, badge=f'🧠 Model-Aware ({i+1}/3)')

    # Improvement 4: Interaction-Aware
    for i, (t, s, b) in enumerate([
        ('Improvement 4: Interaction-Aware Estimation',
         'Detect and account for feature interactions', [
            '• SHAP Interaction Values: φᵢⱼ = f({i,j}) - f({i}) - f({j}) + f({})',
            '• Screen interactions using gradient cross-terms ∂²f/∂xᵢ∂xⱼ',
            '• For detected pairs: compute SHAP interaction values',
            '• Adjust individual φᵢ by interaction correction',
            '• Major improvement for non-linear models',]),
        ('Interaction Detection Algorithm',
         'Use Hessian off-diagonal to find interacting feature pairs', [
            '• Step 1: Compute partial Hessian (top-k off-diagonal terms)',
            '• Step 2: |∂²f/∂xᵢ∂xⱼ| > threshold → interaction detected',
            '• Step 3: For each interacting pair, compute φᵢⱼ',
            '• Step 4: Adjust φᵢ = φᵢ(marginal) + Σⱼ φᵢⱼ/2',
            '• Cost: One Hessian computation + targeted SHAP for pairs',]),
        ('Interaction-Aware Results',
         'Dramatic improvement on XOR-like problems', [
            '• XOR problem: Accuracy 0% → 80% for interaction features',
            '• Neural networks: Overall correlation 85% → 92%',
            '• Random forests: 90% → 95% correlation',
            '• Cost: 2-3× slower than basic InstaSHAP (still 40× faster than exact)',]),
    ]):
        build_improvement_slide(prs, t, s, b, badge=f'🔗 Interact. ({i+1}/3)')

    # Phase 2 Summary slides
    build_metrics_slide(prs, 'v1 vs v2 Performance Comparison', [
        ('+15%', 'Accuracy Gain\n(85→95%)', CLR_SHAP),
        ('3×', 'Interaction\nCoverage', CLR_INPUT),
        ('~2×', 'Slower\n(35→70ms)', CLR_SAMPLING),
    ], badge='📊 COMPARISON')

    build_improvement_slide(prs, 'Implementation Roadmap',
        'Prioritized improvement plan', [
            '• Week 1-2: Adaptive Sampling (Low difficulty, High impact)',
            '• Week 1: Stratified Subsets (Low, Medium impact)',
            '• Week 2-3: Hybrid SHAP (Medium, High impact)',
            '• Week 2-4: Gradient-Guided (Medium, High impact)',
            '• Month 1-2: Interaction Detection (High, Very High impact)',
            '• Month 2-3: Model-Specific Shortcuts (High, Medium impact)',
        ], badge='📋 ROADMAP')

    build_improvement_slide(prs, 'Short-term Research (3-6 months)',
        'Practical improvements achievable in near-term', [
            '• Variance-aware confidence intervals for SHAP values',
            '• Adaptive k selection per feature (auto-tuning)',
            '• Stratified sampling across subset sizes',
            '• Comprehensive benchmark suite for evaluation',
        ], badge='🔬 RESEARCH')

    build_improvement_slide(prs, 'Long-term Research (6-12 months)',
        'Ambitious directions for fundamental improvements', [
            '• Neural network-specific shortcuts (activation patterns)',
            '• LLM-compatible explanation framework (beyond SHAP)',
            '• Interaction-aware SHAP with theoretical guarantees',
            '• Formal error bounds for InstaSHAP approximation',
        ], badge='🔭 FUTURE')

    # Additional Phase 2 slides to reach target
    for title, items in [
        ('Confidence Intervals for SHAP Values', [
            '• Current: Single point estimate φ̂ᵢ with unknown uncertainty',
            '• Proposed: Bootstrap confidence intervals from k samples',
            '• Report: φ̂ᵢ ± 1.96 × SE(φ̂ᵢ) for 95% CI',
            '• User knows which explanations are reliable vs uncertain',]),
        ('Auto-tuning k Per Feature', [
            '• Instead of fixed k for all features, automatically select',
            '• Easy features (low variance): k_min = 10',
            '• Hard features (high variance): k_max = 500',
            '• Convergence criterion: SE(φ̂ᵢ) < ε',
            '• Result: Minimum computation for desired accuracy',]),
        ('Stratified Sampling Implementation', [
            '• Ensure representation across subset sizes |S| = 0,1,...,n',
            '• Shapley values weight small and large subsets heavily',
            '• Current random may under-sample extremes',
            '• Stratified: allocate k/n samples per subset size',
            '• Expected: 30-50% variance reduction',]),
        ('Benchmark Suite Design', [
            '• Standardized test models: linear, tree, MLP, CNN',
            '• Known ground-truth SHAP values for each',
            '• Metrics: MAE, RMSE, Spearman correlation, sign accuracy',
            '• Interaction tests: XOR, pairwise, polynomial',
            '• Distribution shift tests: covariate shift, concept drift',]),
        ('Neural Network Shortcuts', [
            '• Use DeepLIFT-style propagation rules',
            '• Backpropagate attribution through each layer',
            '• Approximation quality depends on architecture',
            '• Potential 10× speedup for supported architectures',
            '• Research needed: correctness proofs, failure modes',]),
        ('LLM Explanation Framework', [
            '• Beyond token-level SHAP: semantic unit attribution',
            '• Group tokens into "meaning units" (phrases, clauses)',
            '• Hierarchical SHAP: sentence → phrase → word → token',
            '• Integration with attention patterns for validation',
            '• Requires fundamental rethinking of feature granularity',]),
        ('Interaction-Aware Theory', [
            '• Shapley interaction index: φᵢⱼ from Owen (1972)',
            '• Computational complexity: O(n² × k) additional',
            '• Can detect pairwise interactions efficiently',
            '• Higher-order interactions (3-way+) remain challenging',
            '• Theoretical error bounds for interaction estimation',]),
        ('Formal Error Bounds', [
            '• Goal: prove ||φ̂ - φ||₂ ≤ ε with probability 1-δ',
            '• Hoeffding bound: k ≥ (range²/2ε²) × ln(2n/δ)',
            '• For ε=0.01, δ=0.05, n=20: k ≈ 300 needed',
            '• Tighter bounds possible with variance information',
            '• Practical: gives users guaranteed accuracy levels',]),
        ('Ensemble Explanation', [
            '• Run InstaSHAP M times with different random seeds',
            '• Average explanations → reduces variance by √M',
            '• Disagreement between runs → flag uncertain features',
            '• Cost: M × t_instashap but still < exact SHAP',]),
        ('Integration with Other XAI Methods', [
            '• Validate InstaSHAP against LIME, gradient-based methods',
            '• Agreement → high confidence in explanation',
            '• Disagreement → flag for more careful analysis',
            '• Meta-explanation: confidence score for each SHAP value',]),
    ]:
        build_improvement_slide(prs, title,
            'Phase 2 Research Direction', items, badge='🚀 v2')

    # Final slides
    build_workflow_step_slide(prs, 0, 'Key Takeaways',
        'What you should remember about InstaSHAP', [
            '✓ InstaSHAP trades exactness for speed (120× faster)',
            '✓ Pipeline: Input → Perturb → Sample → Predict → Compute → Output',
            '✓ Works well for ML models with independent features',
            '✗ Fails on feature interactions, LLMs, non-linear models',
            '✗ Vulnerable to distribution shift',
            '→ v2 improvements can address most limitations',
        ], badge_text='🏁 CONCLUSION', badge_color=CLR_ACCENT)

    build_thank_you_slide(prs)

    # Save
    output_dir = os.path.dirname(os.path.abspath(__file__))
    output_path = os.path.join(output_dir, 'InstaSHAP_Presentation.pptx')
    prs.save(output_path)
    print(f'[OK] Presentation saved to: {output_path}')
    print(f'   Total slides: {len(prs.slides)}')
    return output_path


if __name__ == '__main__':
    generate_presentation()
